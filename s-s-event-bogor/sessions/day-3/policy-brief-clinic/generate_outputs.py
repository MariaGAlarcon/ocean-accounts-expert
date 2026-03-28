#!/usr/bin/env python3
"""
Generate GOAP-styled Word and PowerPoint files for the policy brief clinic.

Outputs:
  - policy-brief-exercise.docx      (exercise worksheet, GOAP Brief template)
  - policy-brief-message-box.docx   (printable Message Box, GOAP Brief template)
  - policy-brief-template.docx      (fill-in brief template, GOAP Policy Brief template)
  - policy-brief-guide.docx         (facilitation guide, GOAP Report template)
  - policy-brief-slides.pptx        (presentation slides, GOAP PPT template)
"""

import re
import sys
from pathlib import Path

# Add paths for imports
BELIZE = Path("/Users/mariaalarcon/GitHub/ocean-accounts-expert/belize-project")
CLINIC = Path("/Users/mariaalarcon/GitHub/ocean-accounts-expert/s-s-event-bogor/sessions/day-2/clinic")
sys.path.insert(0, str(CLINIC))
sys.path.insert(0, str(BELIZE))

from md_to_docx_improved import convert_file

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
HERE = Path(__file__).parent
TEMPLATES = Path("/Users/mariaalarcon/GitHub/ocean-accounts-expert/goap-templates")
REPORT_TEMPLATE = TEMPLATES / "GOAP Word templates" / "GOAP-REPORT TEMPLATE.docx"
BRIEF_TEMPLATE = TEMPLATES / "GOAP Word templates" / "GOAP-BRIEF TEMPLATE.docx"
POLICY_BRIEF_TEMPLATE = TEMPLATES / "GOAP Word templates" / "GOAP-POLICY-BRIEF-TEMPLATE.docx"
PPT_TEMPLATE = TEMPLATES / "GOAP PowerPoint" / "GOAP-PPT-TEMPLATE.pptx"


def pick_word_template(name):
    """Choose the most appropriate Word template for each file."""
    if name == "policy-brief-template":
        # The fill-in template should use the policy brief template if available
        if POLICY_BRIEF_TEMPLATE.exists():
            return POLICY_BRIEF_TEMPLATE
        if BRIEF_TEMPLATE.exists():
            return BRIEF_TEMPLATE
    if name in ("policy-brief-exercise", "policy-brief-message-box"):
        if BRIEF_TEMPLATE.exists():
            return BRIEF_TEMPLATE
    # Default to report template for the guide
    return REPORT_TEMPLATE


# ---------------------------------------------------------------------------
# Word document generation
# ---------------------------------------------------------------------------
def generate_word_docs():
    """Convert markdown files to GOAP-styled Word documents."""
    md_files = {
        "policy-brief-exercise": "Policy Brief Exercise Worksheet",
        "policy-brief-message-box": "COMPASS Message Box -- Ocean Accounts",
        "policy-brief-template": "Policy Brief Template",
        "policy-brief-guide": "Policy Brief Clinic -- Facilitation Guide",
    }

    for name, header in md_files.items():
        md_path = HERE / f"{name}.md"
        docx_path = HERE / f"{name}.docx"
        template = pick_word_template(name)

        if not md_path.exists():
            print(f"  SKIP: {md_path.name} not found")
            continue
        if not template.exists():
            print(f"  SKIP: template {template.name} not found")
            continue

        print(f"Converting {name}.md -> {name}.docx (template: {template.name})")
        convert_file(str(md_path), str(docx_path), str(template), header=header)


# ---------------------------------------------------------------------------
# PowerPoint generation
# ---------------------------------------------------------------------------
def generate_pptx():
    """Generate GOAP-styled PowerPoint from the slides markdown."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    md_path = HERE / "policy-brief-slides.md"
    pptx_path = HERE / "policy-brief-slides.pptx"

    if not md_path.exists():
        print(f"  SKIP: {md_path.name} not found")
        return

    # GOAP colours
    DARK_TEAL = RGBColor(0x0A, 0x54, 0x55)
    GREEN = RGBColor(0x3B, 0x9C, 0x7B)
    BODY = RGBColor(0x30, 0x30, 0x2F)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_MINT = RGBColor(0xD4, 0xEE, 0xE5)

    # Load template or create blank
    if PPT_TEMPLATE.exists():
        prs = Presentation(str(PPT_TEMPLATE))
        print(f"  Using GOAP PPT template: {PPT_TEMPLATE.name}")
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        print("  Using blank presentation (GOAP PPT template not found)")

    # Parse slides from markdown
    text = md_path.read_text(encoding="utf-8")
    # Split on --- slide separators (lines that are just ---)
    raw_slides = re.split(r"\n---\n", text)

    # Detect available layouts
    layouts = {layout.name: layout for layout in prs.slide_layouts}
    # Prefer blank layout for maximum control
    blank_layout = None
    for name in ["Blank", "blank", "Custom Layout"]:
        if name in layouts:
            blank_layout = layouts[name]
            break
    if blank_layout is None:
        # Fall back to last layout (often blank) or first
        blank_layout = prs.slide_layouts[-1]

    # Title slide layout
    title_layout = None
    for name in ["Title Slide", "title", "Title"]:
        if name in layouts:
            title_layout = layouts[name]
            break
    if title_layout is None:
        title_layout = prs.slide_layouts[0]

    # Remove any existing slides from template
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        if rId is None:
            # Try the 'id' attribute
            slide_id_elem = prs.slides._sldIdLst[0]
            # Remove from relationships and slide list
            prs.part.drop_rel(rId) if rId else None
            prs.slides._sldIdLst.remove(slide_id_elem)
        else:
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    def add_slide(layout=None):
        if layout is None:
            layout = blank_layout
        return prs.slides.add_slide(layout)

    def add_textbox(slide, left, top, width, height):
        return slide.shapes.add_textbox(left, top, width, height)

    def set_text(tf, text, font_size=18, color=BODY, bold=False, alignment=PP_ALIGN.LEFT):
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = alignment
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = "Inter"
        return p

    def add_paragraph(tf, text, font_size=16, color=BODY, bold=False, space_before=6,
                      bullet=False):
        p = tf.add_paragraph()
        p.space_before = Pt(space_before)
        if bullet:
            p.level = 0
            # Add bullet character
            text = "  " + text
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = "Inter"
        return p

    # Slide dimensions
    sw = prs.slide_width
    sh = prs.slide_height
    margin = Inches(0.8)
    content_width = sw - 2 * margin
    title_top = Inches(0.6)
    body_top = Inches(1.5)
    body_height = sh - body_top - Inches(0.5)

    for idx, raw in enumerate(raw_slides):
        raw = raw.strip()
        if not raw:
            continue

        lines = raw.split("\n")

        # Extract title (first heading)
        slide_title = ""
        body_lines = []
        found_title = False
        for line in lines:
            m = re.match(r"^#{1,3}\s+(.*)", line)
            if m and not found_title:
                slide_title = m.group(1).strip()
                found_title = True
                continue
            body_lines.append(line)

        # Skip the preamble (metadata before first slide)
        if idx == 0 and not slide_title:
            # This is the front matter, make it the title slide
            for line in lines:
                m = re.match(r"^#\s+(.*)", line)
                if m:
                    slide_title = m.group(1).strip()
                    break
            if not slide_title:
                continue

        # Create slide
        if idx == 0:
            slide = add_slide(title_layout)
            # Try to use placeholder for title
            if slide.placeholders:
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 0:  # title
                        ph.text = slide_title
                        for run in ph.text_frame.paragraphs[0].runs:
                            run.font.name = "Inter"
                            run.font.size = Pt(36)
                            run.font.color.rgb = DARK_TEAL
                        break
                # Add subtitle info
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 1:  # subtitle
                        ph.text = "Ocean Accounting Clinic Pt.2\nDay 3 -- 1 April 2026"
                        for para in ph.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.name = "Inter"
                                run.font.size = Pt(18)
                                run.font.color.rgb = BODY
                        break
            continue

        slide = add_slide(blank_layout)

        # Add title bar with green background
        title_shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0), Inches(0),
            sw, Inches(1.1)
        )
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = GREEN
        title_shape.line.fill.background()
        tf = title_shape.text_frame
        tf.word_wrap = True
        tf.margin_left = margin
        tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = slide_title
        run.font.size = Pt(28)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = "Inter"

        # Process body content
        body_text = "\n".join(body_lines).strip()
        if not body_text:
            continue

        # Add body textbox
        txBox = add_textbox(slide, margin, Inches(1.3), content_width, body_height)
        tf = txBox.text_frame
        tf.word_wrap = True

        # Parse body into elements
        in_table = False
        table_lines = []
        first_para = True

        for line in body_lines:
            stripped = line.strip()

            # Skip empty lines
            if not stripped:
                continue

            # Skip sub-headings in the body (already used as title)
            if stripped.startswith("## Slide"):
                continue

            # Table detection
            if stripped.startswith("|"):
                table_lines.append(stripped)
                in_table = True
                continue
            elif in_table:
                # End of table -- add it as formatted text
                _add_table_to_slide(slide, table_lines, margin, tf, content_width,
                                    DARK_TEAL, GREEN, WHITE, BODY)
                table_lines = []
                in_table = False

            # Code blocks
            if stripped.startswith("```"):
                continue

            # Bullet points
            if stripped.startswith("- ") or stripped.startswith("* "):
                text = stripped[2:]
                text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)  # strip bold markers
                p = add_paragraph(tf, text, font_size=14, color=BODY,
                                  space_before=4, bullet=True)
                first_para = False
                continue

            # Sub-sub headings within slide
            m = re.match(r"^#{2,4}\s+(.*)", stripped)
            if m:
                heading_text = m.group(1).strip()
                heading_text = re.sub(r"\*\*(.+?)\*\*", r"\1", heading_text)
                add_paragraph(tf, heading_text, font_size=18, color=DARK_TEAL,
                              bold=True, space_before=12)
                first_para = False
                continue

            # Regular paragraph
            clean = re.sub(r"\*\*(.+?)\*\*", r"\1", stripped)
            clean = re.sub(r"\*(.+?)\*", r"\1", clean)
            if first_para:
                set_text(tf, clean, font_size=14, color=BODY)
                first_para = False
            else:
                add_paragraph(tf, clean, font_size=14, color=BODY, space_before=6)

        # Handle trailing table
        if table_lines:
            _add_table_to_slide(slide, table_lines, margin, tf, content_width,
                                DARK_TEAL, GREEN, WHITE, BODY)

    # References slide
    slide = add_slide(blank_layout)
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), sw, Inches(1.1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = GREEN
    title_shape.line.fill.background()
    tf = title_shape.text_frame
    tf.word_wrap = True
    tf.margin_left = margin
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "References"
    run.font.size = Pt(28)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "Inter"

    refs_box = add_textbox(slide, margin, Inches(1.5), content_width, Inches(4))
    tf = refs_box.text_frame
    tf.word_wrap = True
    refs = [
        'Baron, N. (2010). Escape from the Ivory Tower: A Guide to Making Your Science Matter. Island Press.',
        'Cairney, P. and Heikkila, T. (2014). A comparison of theories of the policy process. In Theories of the Policy Process. 3rd ed. Westview Press.',
        'Evans, M.C. and Cvitanovic, C. (2018). An introduction to achieving policy impact for early career researchers. Palgrave Communications 4(88).',
        'COMPASS. The Message Box Workbook. Communication Partnership for Science and the Sea.',
    ]
    set_text(tf, refs[0], font_size=12, color=BODY)
    for r in refs[1:]:
        add_paragraph(tf, r, font_size=12, color=BODY, space_before=8)

    prs.save(str(pptx_path))
    print(f"  Created: {pptx_path}")


def _add_table_to_slide(slide, table_lines, margin, text_frame, content_width,
                         dark_teal, green, white, body_color):
    """Add a markdown table to a slide as a PowerPoint table shape."""
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    def split_row(line):
        cells = line.strip().strip("|").split("|")
        return [c.strip() for c in cells]

    if len(table_lines) < 2:
        return

    headers = split_row(table_lines[0])
    rows = []
    for line in table_lines[1:]:
        cells = split_row(line)
        if all(re.match(r"^[-:]+$", c) for c in cells if c):
            continue
        rows.append(cells)

    if not rows:
        return

    num_cols = len(headers)
    num_rows = len(rows) + 1

    # Position table below the title bar area
    table_top = Inches(3.5)

    # Use a reasonable table size
    table_width = min(content_width, Inches(11))
    table_height = Inches(0.35 * num_rows)

    tbl = slide.shapes.add_table(
        num_rows, num_cols,
        margin, table_top,
        table_width, table_height
    ).table

    # Style header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        run.font.name = "Inter"
        run.font.size = Pt(11)
        run.font.color.rgb = white
        run.font.bold = True
        # Green background
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = green

    # Style data rows
    for r_idx, row_data in enumerate(rows):
        for j, val in enumerate(row_data):
            if j >= num_cols:
                continue
            cell = tbl.cell(r_idx + 1, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            # Strip markdown bold markers
            val = re.sub(r"\*\*(.+?)\*\*", r"\1", val)
            run = p.add_run()
            run.text = val
            run.font.name = "Inter"
            run.font.size = Pt(10)
            run.font.color.rgb = body_color


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    print("=" * 60)
    print("Generating policy brief clinic outputs")
    print("=" * 60)
    print()
    print("--- Word documents ---")
    generate_word_docs()
    print()
    print("--- PowerPoint slides ---")
    generate_pptx()
    print()
    print("Done.")
