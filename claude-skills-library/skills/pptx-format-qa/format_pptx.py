#!/usr/bin/env python3
"""
PowerPoint Format QA -- Check and fix presentation formatting.

Checks:
1. Font size minimums (title >= 28pt, body >= 16pt, table >= 12pt)
2. Table readability (flag dense tables)
3. Font consistency (normalize to Arial)
4. GOAP color application

Usage:
    python format_pptx.py path/to/file.pptx
    python format_pptx.py path/to/folder/
    python format_pptx.py path/to/folder/ --recursive

Requires: python-pptx (pip install python-pptx)
"""

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ---- Configuration ----

DEFAULT_CONFIG = {
    "title_min_pt": 28,
    "title_target_pt": 36,
    "body_min_pt": 16,
    "body_target_pt": 18,
    "table_min_pt": 12,
    "table_target_pt": 14,
    "caption_min_pt": 11,
    "absolute_min_pt": 10,
    "font_name": "Arial",
    "max_table_rows": 6,
    "max_table_cols": 5,
    # GOAP colors
    "title_color": RGBColor(0x0A, 0x54, 0x55),  # Teal
    "accent_color": RGBColor(0x3B, 0x9C, 0x7B),  # Green
    "body_color": RGBColor(0x54, 0x53, 0x53),    # Dark gray
    "table_header_bg": RGBColor(0x3B, 0x9C, 0x7B),
    "table_header_text": RGBColor(0xFF, 0xFF, 0xFF),
}


def is_title_shape(shape):
    """Check if a shape is likely a slide title."""
    try:
        if shape.is_placeholder:
            idx = shape.placeholder_format.idx
            return idx in (0, 1)  # Title or subtitle placeholder
    except (ValueError, AttributeError):
        pass
    # Fallback: check if text is short and has large font
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if len(text) < 80:  # Short text likely a title
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size / 12700 >= 24:
                        return True
    return False


def process_file(filepath, config=None, fix=True):
    """Check and optionally fix a PowerPoint file."""
    cfg = config or DEFAULT_CONFIG
    prs = Presentation(str(filepath))
    issues = []
    fixes_applied = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:

            # ---- Text frames ----
            if shape.has_text_frame:
                is_title = is_title_shape(shape)

                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue

                        current_size = run.font.size
                        if current_size is None:
                            continue

                        size_pt = current_size / 12700  # EMU to points

                        # Determine minimum based on context
                        if is_title:
                            min_size = cfg["title_min_pt"]
                            target = cfg["title_target_pt"]
                            context = "title"
                        else:
                            min_size = cfg["body_min_pt"]
                            target = cfg["body_target_pt"]
                            context = "body"

                        # Check and fix font size
                        if size_pt < min_size:
                            issues.append(
                                f"Slide {slide_num}: {context} text "
                                f'"{run.text[:30]}..." is {size_pt:.0f}pt '
                                f"(min {min_size}pt)"
                            )
                            if fix:
                                run.font.size = Pt(target)
                                fixes_applied += 1

                        # Check and fix font name
                        if fix and run.font.name != cfg["font_name"]:
                            run.font.name = cfg["font_name"]

            # ---- Tables ----
            if shape.has_table:
                table = shape.table
                rows = len(table.rows)
                cols = len(table.columns)

                if rows > cfg["max_table_rows"]:
                    issues.append(
                        f"Slide {slide_num}: table has {rows} rows "
                        f"(max recommended {cfg['max_table_rows']})"
                    )

                if cols > cfg["max_table_cols"]:
                    issues.append(
                        f"Slide {slide_num}: table has {cols} columns "
                        f"(max recommended {cfg['max_table_cols']})"
                    )

                for row_idx, row in enumerate(table.rows):
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if not run.text.strip():
                                    continue

                                current_size = run.font.size
                                if current_size is None:
                                    continue

                                size_pt = current_size / 12700

                                if size_pt < cfg["table_min_pt"]:
                                    issues.append(
                                        f"Slide {slide_num}: table text "
                                        f'"{run.text[:20]}..." is '
                                        f"{size_pt:.0f}pt "
                                        f"(min {cfg['table_min_pt']}pt)"
                                    )
                                    if fix:
                                        run.font.size = Pt(cfg["table_target_pt"])
                                        fixes_applied += 1

                                # Fix font
                                if fix and run.font.name != cfg["font_name"]:
                                    run.font.name = cfg["font_name"]

                                # Apply GOAP colors to header row
                                if fix and row_idx == 0:
                                    run.font.bold = True
                                    run.font.color.rgb = cfg["table_header_text"]

                        # Header row background
                        if fix and row_idx == 0:
                            from pptx.oxml.ns import qn
                            tcPr = cell._tc.get_or_add_tcPr()
                            solidFill = tcPr.find(qn("a:solidFill"))
                            if solidFill is None:
                                from lxml import etree
                                solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
                                srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
                                srgbClr.set("val", "3B9C7B")

    if fix and (issues or fixes_applied):
        prs.save(str(filepath))

    return issues, fixes_applied


def process_path(target, recursive=False, config=None, fix=True):
    """Process a file or folder."""
    target = Path(target)
    total_files = 0
    total_issues = 0
    total_fixes = 0

    if target.is_file() and target.suffix == ".pptx":
        files = [target]
    elif target.is_dir():
        pattern = "**/*.pptx" if recursive else "*.pptx"
        files = sorted(target.glob(pattern))
    else:
        print(f"Error: {target} is not a .pptx file or directory")
        return

    for filepath in files:
        if filepath.name.startswith("~$"):
            continue
        total_files += 1
        issues, fixes = process_file(filepath, config, fix)
        total_issues += len(issues)
        total_fixes += fixes

        if issues:
            print(f"\n  {filepath.name}: {len(issues)} issues, {fixes} fixes")
            for issue in issues[:10]:  # Show first 10
                print(f"    - {issue}")
            if len(issues) > 10:
                print(f"    ... and {len(issues) - 10} more")
        else:
            print(f"  [CLEAN] {filepath.name}")

    print(f"\nTotal: {total_files} files, {total_issues} issues found, {total_fixes} fixes applied")


def main():
    parser = argparse.ArgumentParser(
        description="PowerPoint Format QA: check and fix presentation formatting"
    )
    parser.add_argument("path", help="Path to .pptx file or folder")
    parser.add_argument("--recursive", "-r", action="store_true")
    parser.add_argument("--check-only", action="store_true",
                        help="Report issues without fixing them")
    parser.add_argument("--title-min", type=int, default=28,
                        help="Minimum title font size in pt (default: 28)")
    parser.add_argument("--body-min", type=int, default=16,
                        help="Minimum body font size in pt (default: 16)")
    parser.add_argument("--table-min", type=int, default=12,
                        help="Minimum table font size in pt (default: 12)")

    args = parser.parse_args()

    config = DEFAULT_CONFIG.copy()
    config["title_min_pt"] = args.title_min
    config["body_min_pt"] = args.body_min
    config["table_min_pt"] = args.table_min

    print("PowerPoint Format QA")
    print(f"Target: {args.path}")
    print(f"Mode: {'Check only' if args.check_only else 'Check and fix'}")
    print(f"{'='*60}")

    process_path(args.path, args.recursive, config, fix=not args.check_only)


if __name__ == "__main__":
    main()
